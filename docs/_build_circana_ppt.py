"""Build the Circana AI-for-BI pitch deck (PPTX) mirroring the HTML structure."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- Circana palette (existing) --------------------------------------------
MAGENTA = RGBColor(0xC8, 0x38, 0x9B)
PURPLE  = RGBColor(0x6B, 0x2C, 0x8E)
DEEP    = RGBColor(0x2A, 0x0F, 0x40)
INK     = RGBColor(0x15, 0x15, 0x1A)
INK2    = RGBColor(0x4A, 0x4A, 0x55)
LINE    = RGBColor(0xE6, 0xE1, 0xEE)
SOFT    = RGBColor(0xF7, 0xF4, 0xFB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
WARM    = RGBColor(0xFF, 0x6A, 0x3D)

# ---- Operativo palette (warm yellow-orange origami + lime green accent) ----
OP_YELLOW   = RGBColor(0xFA, 0xB6, 0x17)
OP_ORANGE   = RGBColor(0xFF, 0x7A, 0x3D)
OP_GREEN    = RGBColor(0x6F, 0xCF, 0x3C)
OP_DEEP     = RGBColor(0x1B, 0x1B, 0x22)
OP_LINE     = RGBColor(0xE8, 0xE4, 0xDA)
OP_SOFT     = RGBColor(0xFE, 0xF8, 0xEB)

# ---- GSK palette (from tailwind config) ------------------------------------
GSK_SUNSET    = RGBColor(0xFF, 0xAA, 0x33)
GSK_ORANGE    = RGBColor(0xF3, 0x66, 0x33)
GSK_DEEP      = RGBColor(0xE0, 0x3C, 0x31)
GSK_CHARCOAL  = RGBColor(0x2D, 0x2D, 0x2D)
GSK_SLATE     = RGBColor(0x5A, 0x5A, 0x5A)
GSK_CLOUD     = RGBColor(0xF5, 0xF5, 0xF5)
GSK_LINE      = RGBColor(0xE8, 0xE8, 0xE8)
GSK_TEAL      = RGBColor(0x00, 0xA3, 0x9E)

ROOT = Path(__file__).resolve().parent.parent
LOGO = ROOT / "circana-logo.png"
OUT  = ROOT / "docs" / "circana_ai_for_bi_pitch.pptx"

# External logo files (other Accenture pitches included in the portfolio deck)
OPERATIVO_LOGO = Path(r"C:\Users\sudheer.chinthala\Downloads\campaignCockpit\operativo-logo.png")
GSK_LOGO       = Path(r"C:\Users\sudheer.chinthala\Downloads\NovoNordisk\GSK-logo.png")

# Widescreen 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ---- helpers ----------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill, line=None, corner=0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = corner
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_richtext(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, font="Calibri"):
    """runs = list of (text, size, bold, color) or (text, size, bold, color, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for spec in runs:
        if len(spec) == 4:
            text, size, bold, color = spec
            italic = False
        else:
            text, size, bold, color, italic = spec
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def gradient_band(slide, x, y, w, h):
    """Approximate the deep->purple->magenta gradient by stacking rectangles."""
    # solid base
    base = add_rect(slide, x, y, w, h, DEEP)
    # Use a real gradient fill via XML for the band
    spPr = base.fill._xPr.find(qn("a:solidFill"))
    if spPr is not None:
        base.fill._xPr.remove(spPr)
    gradFill = etree.SubElement(base.fill._xPr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    stops = [(0, DEEP), (55000, PURPLE), (100000, MAGENTA)]
    for pos, col in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", "2700000")  # ~45 degrees
    lin.set("scaled", "1")
    return base


def slide_header(slide, num, title, sub=None):
    """Standard inner-slide header: section number chip + title + thin accent line."""
    # top accent stripe
    add_rect(slide, 0, 0, SW, Inches(0.18), MAGENTA)
    # number chip
    chip = add_rounded(slide, Inches(0.5), Inches(0.42), Inches(0.7), Inches(0.45),
                       PURPLE, corner=0.3)
    add_text(slide, Inches(0.5), Inches(0.42), Inches(0.7), Inches(0.45),
             num, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_text(slide, Inches(1.35), Inches(0.38), Inches(11), Inches(0.55),
             title, size=24, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    # accent line under title
    add_rect(slide, Inches(0.5), Inches(1.0), Inches(12.3), Emu(9525), LINE)
    if sub:
        add_text(slide, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.45),
                 sub, size=12, color=INK2)


def footer_bar(slide, page_num, total):
    add_rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), SOFT)
    add_text(slide, Inches(0.5), SH - Inches(0.30), Inches(8), Inches(0.28),
             "AI for BI  |  Market Insights Acceleration for Circana",
             size=9, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(2), SH - Inches(0.30), Inches(1.5), Inches(0.28),
             f"{page_num} / {total}",
             size=9, color=INK2, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ---- Slide 1: Title ---------------------------------------------------------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    gradient_band(s, 0, 0, SW, SH)

    # decorative circle (subtle)
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              SW - Inches(3.5), -Inches(2.5),
                              Inches(6), Inches(6))
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    deco.fill.transparency = 0  # python-pptx ignores this; left for clarity
    deco.line.fill.background()
    # transparency hack via XML
    solidFill = deco.fill._xPr.find(qn("a:solidFill"))
    srgb = solidFill.find(qn("a:srgbClr"))
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", "10000")  # ~10%

    # Logo on white plate, top-left
    plate = add_rounded(s, Inches(0.6), Inches(0.5), Inches(2.2), Inches(0.85),
                        WHITE, corner=0.18)
    s.shapes.add_picture(str(LOGO), Inches(0.85), Inches(0.6),
                         height=Inches(0.65))

    # right-top tag
    add_text(s, SW - Inches(5.6), Inches(0.7), Inches(5), Inches(0.4),
             "AI FOR BI   |   MARKET INSIGHTS PITCH BRIEF",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # Title block
    add_text(s, Inches(0.6), Inches(2.0), Inches(11.5), Inches(1.4),
             "From Dashboards to Decisions",
             size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(2.85), Inches(11.5), Inches(1.0),
             "Conversational, Agentic AI for Circana's CPG & GM Market Insights",
             size=24, bold=False, color=WHITE)

    # Subtitle
    add_text(s, Inches(0.6), Inches(4.0), Inches(11.5), Inches(1.1),
             ("A three-layer AI-for-BI reference — natural-language analytics, a "
              "programmable policy engine, and human-in-the-loop agents — purpose-built "
              "for FMCG, CPG, and General Merchandise insights teams, designed to compound "
              "the value of Circana's Liquid Data® and Liquid AI™."),
             size=14, color=WHITE)

    # Stat strip
    stats = [
        ("$5.8T", "Global Consumer Spend tracked by Circana"),
        ("42M",   "CPG & GM items actively measured"),
        ("477K",  "Stores in the panel"),
        ("26",    "Industries — CPG, GM, Beauty, Durables, Tech"),
    ]
    x0 = Inches(0.6)
    w  = (SW - Inches(1.2)) / 4
    y  = Inches(5.6)
    for i, (n, l) in enumerate(stats):
        bx = x0 + i * w + Inches(0.05)
        bw = w - Inches(0.1)
        card = add_rounded(s, bx, y, bw, Inches(1.3), DEEP, corner=0.1)
        # translucent overlay via XML
        sf = card.fill._xPr.find(qn("a:solidFill"))
        srgb = sf.find(qn("a:srgbClr"))
        alpha = etree.SubElement(srgb, qn("a:alpha"))
        alpha.set("val", "35000")
        add_text(s, bx, y + Inches(0.15), bw, Inches(0.55),
                 n, size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bx + Inches(0.1), y + Inches(0.7), bw - Inches(0.2),
                 Inches(0.55), l, size=10, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # bottom meta line
    add_text(s, Inches(0.6), SH - Inches(0.5), Inches(11.5), Inches(0.3),
             "Pitch Brief  |  May 2026",
             size=10, color=WHITE)


# ---- Slide 2: Executive Thesis ----------------------------------------------
def slide_thesis(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "01", "Executive Thesis",
                 ("Circana sits on the largest commercial dataset in CPG and General "
                  "Merchandise. The unsolved problem is the last mile — turning that "
                  "data into a brand, category, or customer-team manager's next action "
                  "in seconds, with auditable, defensible insight."))
    # Three cards
    cards = [
        ("The market gap",
         "CPG and GM brand teams generate enormous POS, panel, media, and price data "
         "but lack the tools to extract insight at the speed decisions require. The "
         "average time from a share-shift signal to a management decision is 2–3 weeks."),
        ("What we offer",
         "A working three-layer reference architecture (NL chatbot → rules engine "
         "→ 7 agents) plus a delivery playbook for standing it up on a new commercial "
         "domain in under 12 weeks — with the audit trail, citation, and governance "
         "an enterprise CPG client expects."),
        ("Why this wins now",
         "It is complementary to Liquid AI, not competitive. Liquid AI is the platform; "
         "we accelerate the persona-shaped surfaces (CPG Brand, Category, Retailer, "
         "Innovation, Customer Team) that a horizontal platform will always need on top."),
    ]
    cw = (SW - Inches(1.4)) / 3
    cy = Inches(1.9)
    ch = Inches(4.6)
    for i, (h, body) in enumerate(cards):
        cx = Inches(0.5) + i * (cw + Inches(0.2))
        add_rounded(s, cx, cy, cw, ch, WHITE, line=LINE, corner=0.05)
        # left accent
        add_rect(s, cx, cy, Inches(0.08), ch, MAGENTA)
        add_text(s, cx + Inches(0.25), cy + Inches(0.25), cw - Inches(0.4),
                 Inches(0.5), h, size=15, bold=True, color=PURPLE)
        add_text(s, cx + Inches(0.25), cy + Inches(0.85), cw - Inches(0.4),
                 ch - Inches(1.1), body, size=12, color=INK2)
    footer_bar(s, idx, total)


# ---- Slide 3: The Asset (3-layer) -------------------------------------------
def slide_asset(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "02", "The Asset — A Three-Layer AI-for-BI Reference Tuned for CPG & GM",
                 ("A working reference across all three layers, configurable for any "
                  "commercial domain Circana tracks — food & beverage, beauty, healthcare, "
                  "durables, tech, fashion, automotive aftermarket, foodservice."))

    # Stacked layer blocks (left 2/3)
    layers = [
        ("LAYER 1  ·  Conversational BI for Market Insights",
         "Intent engine (LLM) → query library against a CPG/GM semantic ontology "
         "(share, velocity, ACV, distribution, price-gap, promo lift, basket affinity) "
         "→ auto-chosen visualization → AI-written narrative. Multi-turn memory; "
         "CPG-native synonym resolution; precise time-expression handling."),
        ("LAYER 2  ·  Policy / Rules-as-a-Service",
         "Programmable rules engine continuously evaluating live data against codified "
         "commercial policy — JBP commitments, retailer contract thresholds, "
         "distribution targets, MAP / MSRP guardrails, brand-health watchlists, "
         "promo-ROI floors, OSA service levels. Violations flow as ranked alerts with "
         "full evidence chain."),
        ("LAYER 3  ·  Agentic Actions for Insights at Scale",
         "Autonomous agents that monitor signals, draft artifacts, propose actions with "
         "human approval gates — Share-Movement Watcher, MBR / Category-Review "
         "Auto-Drafter, Root-Cause Investigator, Innovation / Pacesetter Watcher, "
         "Retailer-Readiness Agent, Insight-Distribution Agent, Pricing & Promo "
         "Investigator."),
    ]
    lx = Inches(0.5)
    lw = Inches(8.3)
    ly = Inches(1.85)
    lh = Inches(1.55)
    for i, (h, body) in enumerate(layers):
        y = ly + i * (lh + Inches(0.12))
        add_rect(s, lx, y, Inches(0.08), lh, MAGENTA)
        add_text(s, lx + Inches(0.2), y + Inches(0.08), lw - Inches(0.3),
                 Inches(0.35), h, size=12, bold=True, color=DEEP)
        add_text(s, lx + Inches(0.2), y + Inches(0.45), lw - Inches(0.3),
                 lh - Inches(0.5), body, size=11, color=INK2)

    # Right panel: 3 soft cards stacked
    right_cards = [
        ("Tech stack (transferable)",
         "FastAPI · LLM-agnostic (Groq, Azure OpenAI, Bedrock, on-prem Llama) "
         "· DuckDB / Snowflake / BigQuery · Malloy semantic layer "
         "· Vega-Lite · immutable audit log."),
        ("Already built",
         "13 chart types, multi-domain query library, follow-up resolution, session "
         "memory, agent runtime with tool-calling, policy DSL, narrative generator."),
        ("Configured for Circana",
         "Ontology (CPG/GM measures & dimensions), query library, and policy pack. "
         "Everything else is the same engine — 12-week vertical pilot is realistic."),
    ]
    rx = Inches(9.0)
    rw = Inches(3.85)
    ry = Inches(1.85)
    rh = Inches(1.55)
    for i, (h, body) in enumerate(right_cards):
        y = ry + i * (rh + Inches(0.12))
        add_rounded(s, rx, y, rw, rh, SOFT, line=LINE, corner=0.05)
        add_text(s, rx + Inches(0.15), y + Inches(0.1), rw - Inches(0.3),
                 Inches(0.35), h, size=11, bold=True, color=PURPLE)
        add_text(s, rx + Inches(0.15), y + Inches(0.45), rw - Inches(0.3),
                 rh - Inches(0.5), body, size=10, color=INK2)

    footer_bar(s, idx, total)


# ---- Slide 4: Workflow mapping table ----------------------------------------
def slide_mapping(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "03", "The Insight-Workflow Mapping",
                 ("Every recurring CPG / GM analyst workflow has a direct AI-for-BI "
                  "capability that replaces it — measured against today's analyst-led baseline."))
    rows = [
        ("Analyst Workflow (CPG / GM)", "Today's Baseline", "With AI for BI"),
        ("Weekly share & velocity review", "4–8 hours of analyst pulls + slides",
         "30-second NL query, auto-narrative, auto-chart"),
        ("Monthly MBR / Category Review drafting", "8–20 hours per report",
         "Agent-drafted deck with cited sources; analyst edits"),
        ("Share-loss / share-gain root cause", "1–2 weeks across multiple cuts",
         "1–2 hours via multi-hop agent with ranked hypotheses"),
        ("Promo ROI / lift deep-dive", "Bespoke analyst engagement, 3–5 days",
         "Conversational drill: lift, halo, cannibalization in minutes"),
        ("Innovation / Pacesetter scan", "Annual report cycle",
         "Continuous agent flagging Year-1 launches 6 months earlier"),
        ("Retailer JBP / Top-to-Top prep", "5 days, multi-team",
         "Readiness Agent produces a single brief, 1 day"),
        ("Distribution / OSA gap monitoring", "Discovered at monthly review",
         "Real-time policy alert with evidence chain"),
        ("Macro-trend ripple (GLP-1, weather, BNPL)", "Thought-leadership reports, lagging",
         "Continuous agent surfacing impacted categories & SKUs"),
        ("Ad-hoc executive question turnaround", "2–4 hours per question",
         "Under 60 seconds, with follow-up suggestions"),
    ]
    cols_w = [Inches(4.3), Inches(4.0), Inches(4.4)]
    cols_x = [Inches(0.5), Inches(0.5) + cols_w[0], Inches(0.5) + cols_w[0] + cols_w[1]]
    row_h = Inches(0.42)
    y0 = Inches(1.85)
    # header row
    for i, txt in enumerate(rows[0]):
        add_rect(s, cols_x[i], y0, cols_w[i], row_h, SOFT,
                 line=LINE if i > 0 else None)
        add_text(s, cols_x[i] + Inches(0.12), y0, cols_w[i] - Inches(0.2),
                 row_h, txt.upper(), size=10, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)
    # data rows
    for r, row in enumerate(rows[1:], start=1):
        y = y0 + r * row_h
        bg = WHITE if r % 2 == 1 else SOFT
        for i, txt in enumerate(row):
            add_rect(s, cols_x[i], y, cols_w[i], row_h, bg)
            bold = (i == 0)
            color = INK if i == 0 or i == 2 else INK2
            add_text(s, cols_x[i] + Inches(0.12), y, cols_w[i] - Inches(0.2),
                     row_h, txt, size=10, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
    # bottom line
    add_rect(s, Inches(0.5), y0 + len(rows) * row_h, sum(cols_w, Inches(0)),
             Emu(9525), LINE)

    footer_bar(s, idx, total)


# ---- Slide 5: Top Use Cases (1-3) -------------------------------------------
def slide_usecases_a(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "04", "Top Use Cases to Pitch  (1 of 2)",
                 "Direct configurations of the existing three-layer reference. Ranked for client demand × build feasibility.")
    cases = [
        ("1", "\"Ask Liquid Data\" — Conversational Share, Velocity & Distribution",
         "Brand managers, category leaders, and customer-team analysts ask in plain English: "
         "\"Which categories are losing share in Kroger over the last 8 weeks and why?\" "
         "— answered in seconds with a chart, narrative, and three suggested follow-ups. "
         "Spans CPG (food, beverage, beauty, HBA), GM (toys, home, apparel, tech), and Foodservice.",
         [("Buyer", "VP Insights / Category at CPG & GM"),
          ("Lift", "Low · ontology + query-library swap"),
          ("KPI",  "Insights / analyst / day ×5–10")]),
        ("2", "Auto-Drafted MBRs, Category & Quarterly Business Reviews",
         "Manufacturers produce hundreds of MBRs and category reviews per month, each 6–20 "
         "analyst hours. An agent assembles share, distribution, price, promo, innovation, "
         "and panel slides with cited Circana sources — analyst becomes the editor. "
         "Productizable as a billable Circana service tier.",
         [("Buyer", "Custom Analytics & Client Services"),
          ("Lift", "Medium · templates + orchestration"),
          ("KPI",  "Report cycle 5 days → same-day")]),
        ("3", "\"Why Is Share Moving?\" Root-Cause Agent  —  Force-Multiplier on Complete Why™",
         "Multi-hop reasoning across POS, panel, distribution, price-gap, promo, media spend, "
         "retail media, weather, and macro signals — producing a ranked hypothesis list "
         "with evidence. Direct accelerator on Circana's Complete Why™.",
         [("Buyer", "Brand Director, RGM, Category Captain"),
          ("Lift", "Med-High · Liquid Data tool wrappers"),
          ("KPI",  "Time-to-root-cause 2 wks → 2 hrs")]),
    ]
    _draw_usecases(s, cases, top=Inches(1.85))
    footer_bar(s, idx, total)


def slide_usecases_b(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "04", "Top Use Cases to Pitch  (2 of 2)",
                 "Use cases 4 and 5 — innovation and customer-engagement focused.")
    cases = [
        ("4", "Continuous Innovation & Pacesetter Watch Agent (CPG + GM)",
         "Circana's 2025 New Product Pacesetters drove $6.2B in Year-1 sales. Convert that "
         "annual signal into a continuous agent: scan Year-1 trajectories of every launch, "
         "flag emerging Pacesetters 6 months earlier, identify stalled launches, write a "
         "one-page brief per candidate. Spans CPG and Durables / GM.",
         [("Buyer", "Innovation Heads at CPG & GM; Circana Innovation"),
          ("Lift", "Medium · policy pack + watcher agent"),
          ("KPI",  "Pacesetter call lead-time +6 months")]),
        ("5", "Retailer / Top-to-Top JBP Readiness Agent",
         "90 days before a JBP with Walmart, Kroger, Target, Costco, or Amazon, the agent "
         "assembles category share, distribution gaps, promo effectiveness, OSA, competitor "
         "moves, and recommended asks into a single readiness brief with a score. The "
         "indispensable weekly artifact for any Customer Team Lead. Adjacent to Market Share Drivers.",
         [("Buyer", "Customer Team Leads at top-100 CPG & GM"),
          ("Lift", "Medium · readiness scoring framework"),
          ("KPI",  "Top-to-top prep 5 days → 1 day")]),
    ]
    _draw_usecases(s, cases, top=Inches(1.85), card_h=Inches(2.4))
    footer_bar(s, idx, total)


def _draw_usecases(s, cases, top, card_h=Inches(1.6)):
    x = Inches(0.5)
    w = Inches(12.3)
    y = top
    for num, title, body, kvs in cases:
        add_rounded(s, x, y, w, card_h, WHITE, line=LINE, corner=0.03)
        # badge
        bx = x + Inches(0.2)
        by = y + Inches(0.2)
        add_rounded(s, bx, by, Inches(0.7), Inches(0.7), PURPLE, corner=0.2)
        add_text(s, bx, by, Inches(0.7), Inches(0.7), num,
                 size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # title
        add_text(s, x + Inches(1.05), y + Inches(0.15), w - Inches(1.3),
                 Inches(0.5), title, size=14, bold=True, color=DEEP)
        # body
        add_text(s, x + Inches(1.05), y + Inches(0.6), w - Inches(1.3),
                 card_h - Inches(1.15), body, size=10.5, color=INK2)
        # kvs row
        kv_y = y + card_h - Inches(0.45)
        kv_x = x + Inches(1.05)
        per = (w - Inches(1.3)) / len(kvs)
        for i, (k, v) in enumerate(kvs):
            add_richtext(s, kv_x + i * per, kv_y, per - Inches(0.1), Inches(0.4),
                         [(k.upper() + "  ", 9, True, PURPLE),
                          (v, 10, False, INK)],
                         anchor=MSO_ANCHOR.MIDDLE)
        y = y + card_h + Inches(0.12)


# ---- Slide 6: Liquid AI positioning -----------------------------------------
def slide_position(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "05", "Positioning vs. Liquid AI — Complement, Don't Compete",
                 ("Liquid AI is Circana's horizontal GenAI platform. The white space is the "
                  "verticalized, persona-shaped analyst surfaces on top — a category a "
                  "horizontal platform cannot fully ship itself."))
    # Two columns
    pairs = [
        ("Where Liquid AI sits",
         "Data plane, retrieval, governance, LLM orchestration, security boundary — "
         "the platform plumbing every persona surface will build on. Strategic to Circana; "
         "not the place a partner should try to substitute."),
        ("Where we sit on top",
         "Domain-specific intent engines, policy packs (JBP, MAP, distribution), ready-built "
         "agents per persona (Brand, Category, Retailer, Innovation, Customer Team), and the "
         "deliverable-quality narrative + audit layer that lets a Brand Director defend an "
         "insight to their CFO."),
    ]
    cw = (SW - Inches(1.2)) / 2
    cy = Inches(1.85)
    ch = Inches(3.2)
    for i, (h, body) in enumerate(pairs):
        cx = Inches(0.5) + i * (cw + Inches(0.2))
        add_rounded(s, cx, cy, cw, ch, WHITE, line=LINE, corner=0.04)
        add_rect(s, cx, cy, Inches(0.08), ch, PURPLE)
        add_text(s, cx + Inches(0.25), cy + Inches(0.25), cw - Inches(0.4),
                 Inches(0.5), h, size=16, bold=True, color=PURPLE)
        add_text(s, cx + Inches(0.25), cy + Inches(0.85), cw - Inches(0.4),
                 ch - Inches(1.1), body, size=12.5, color=INK2)

    # Callout box (gradient)
    cy2 = Inches(5.3)
    ch2 = Inches(1.5)
    cx2 = Inches(0.5)
    cw2 = SW - Inches(1)
    cal = add_rounded(s, cx2, cy2, cw2, ch2, DEEP, corner=0.04)
    # gradient fill on callout
    sf = cal.fill._xPr.find(qn("a:solidFill"))
    cal.fill._xPr.remove(sf)
    g = etree.SubElement(cal.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in [(0, DEEP), (100000, PURPLE)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "0"); lin.set("scaled", "1")
    add_text(s, cx2 + Inches(0.3), cy2 + Inches(0.18), cw2 - Inches(0.6),
             Inches(0.4), "The one-sentence pitch to Circana",
             size=13, bold=True, color=WHITE)
    add_text(s, cx2 + Inches(0.3), cy2 + Inches(0.6), cw2 - Inches(0.6),
             ch2 - Inches(0.7),
             ("\"You've built the platform — let us build the personas. We bring a proven "
              "AI-for-BI reference, port it onto Liquid Data and Liquid AI, and stand up two "
              "verticalized analyst experiences inside one quarter — defensible insights, "
              "auditable narratives, agent-driven artifacts, ready to monetize.\""),
             size=12, color=WHITE)
    footer_bar(s, idx, total)


# ---- Slide 7: Engagement path -----------------------------------------------
def slide_engagement(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "06", "Proposed Engagement Path",
                 "Land with a paid 12-week proof-of-value on one use case; expand into a multi-quarter build-out.")
    phases = [
        ("Weeks 0–2", "Discover",
         "Co-select use case #1 (recommended: Ask Liquid Data + auto-MBR). Map the "
         "Liquid Data semantic layer to our ontology shape. Identify 5 reference "
         "clients. Workshops, data access, persona interviews."),
        ("Weeks 3–8", "Build PoV",
         "Configure intent engine on the Circana ontology; stand up the query library "
         "for the chosen domain; integrate Vega rendering into a Circana-branded "
         "surface; codify 5–10 policy rules. Working demo on real client data by week 8."),
        ("Weeks 9–12", "Pilot & Measure",
         "Pilot with 2 client design partners; instrument time-to-insight, "
         "report-cycle-time, adoption; deliver business case for scale-up; agree "
         "the commercial model (per-seat, per-report, or platform tier)."),
        ("Q2 onward", "Scale",
         "Add use cases 2–5 in parallel tracks; deepen Liquid AI integration; "
         "productize policy packs per retailer / per category; co-sell into Circana's "
         "top-100 accounts as a value-add tier."),
    ]
    x = Inches(0.5)
    w = SW - Inches(1)
    y = Inches(1.85)
    ph_h = Inches(0.85)
    for label, sub, body in phases:
        add_rounded(s, x, y, w, ph_h, WHITE, line=LINE, corner=0.03)
        # label block
        add_rect(s, x, y, Inches(1.55), ph_h, SOFT)
        add_text(s, x + Inches(0.1), y + Inches(0.08), Inches(1.4),
                 Inches(0.3), label, size=11, bold=True, color=MAGENTA,
                 anchor=MSO_ANCHOR.TOP)
        add_text(s, x + Inches(0.1), y + Inches(0.4), Inches(1.4),
                 Inches(0.4), sub, size=13, bold=True, color=DEEP,
                 anchor=MSO_ANCHOR.TOP)
        add_text(s, x + Inches(1.7), y + Inches(0.12), w - Inches(1.85),
                 ph_h - Inches(0.2), body, size=11, color=INK2,
                 anchor=MSO_ANCHOR.MIDDLE)
        y = y + ph_h + Inches(0.12)

    # De-risk callout
    cy = Inches(5.8)
    ch = Inches(1.1)
    cal = add_rounded(s, x, cy, w, ch, DEEP, corner=0.04)
    sf = cal.fill._xPr.find(qn("a:solidFill"))
    cal.fill._xPr.remove(sf)
    g = etree.SubElement(cal.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in [(0, DEEP), (100000, MAGENTA)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "0"); lin.set("scaled", "1")
    add_text(s, x + Inches(0.3), cy + Inches(0.12), w - Inches(0.6),
             Inches(0.4), "Why this de-risks the deal for Circana",
             size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.3), cy + Inches(0.45), w - Inches(0.6),
             ch - Inches(0.55),
             ("A working three-layer reference exists today — not a slideware demo. "
              "Architecture is LLM-agnostic, runs in their cloud, and was originally built "
              "under one of the most demanding governance regimes in any industry. Circana "
              "is buying speed-to-market and risk transfer, not unproven IP."),
             size=10.5, color=WHITE)
    footer_bar(s, idx, total)


# ---- Slide 8: Next steps + closing ------------------------------------------
def slide_next(idx, total):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "07", "Recommended Next Steps")
    steps = [
        ("Align the offering",
         "Confirm internal naming, commercial model, and the lead practice. Resolve "
         "IP-clearance for the reference assets so they can be demonstrated externally."),
        ("Targeted Circana intro",
         "Request a 45-minute working session with Liquid AI product leadership and one "
         "Client Services VP. Lead with a use case #1 demo on a public CPG dataset (or "
         "synthetic panel) reskinned in Circana branding."),
        ("Reference-client lining",
         "Identify 2 CPG / GM clients who would co-pilot. Their participation de-risks "
         "the Circana commercial conversation and makes the PoV self-funding."),
    ]
    cw = (SW - Inches(1.4)) / 3
    cy = Inches(1.85)
    ch = Inches(3.5)
    for i, (h, body) in enumerate(steps):
        cx = Inches(0.5) + i * (cw + Inches(0.2))
        add_rounded(s, cx, cy, cw, ch, WHITE, line=LINE, corner=0.04)
        # number badge
        bx = cx + Inches(0.25); by = cy + Inches(0.25)
        add_rounded(s, bx, by, Inches(0.55), Inches(0.55), MAGENTA, corner=0.3)
        add_text(s, bx, by, Inches(0.55), Inches(0.55), str(i + 1),
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.95), cy + Inches(0.32), cw - Inches(1.15),
                 Inches(0.45), h, size=14, bold=True, color=DEEP)
        add_text(s, cx + Inches(0.25), cy + Inches(1.1), cw - Inches(0.45),
                 ch - Inches(1.25), body, size=12, color=INK2)

    # Closing band
    band_y = Inches(5.7)
    band = add_rounded(s, Inches(0.5), band_y, SW - Inches(1), Inches(1.2),
                       DEEP, corner=0.04)
    sf = band.fill._xPr.find(qn("a:solidFill"))
    band.fill._xPr.remove(sf)
    g = etree.SubElement(band.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in [(0, DEEP), (55000, PURPLE), (100000, MAGENTA)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "0"); lin.set("scaled", "1")
    add_text(s, Inches(0.75), band_y + Inches(0.18), SW - Inches(1.5),
             Inches(0.5), "Ready when you are.",
             size=20, bold=True, color=WHITE)
    add_text(s, Inches(0.75), band_y + Inches(0.62), SW - Inches(1.5),
             Inches(0.5),
             "AI for BI  ·  Market Insights Acceleration for Circana  ·  Pitch Brief, May 2026",
             size=11, color=WHITE)
    footer_bar(s, idx, total)


# ---- Build ------------------------------------------------------------------
# ============================================================================
# OPERATIVO  ·  CONGRESS TWIN  (slides 10-11)
# ============================================================================
# Brand: warm yellow-orange origami mark, lime-green accent dot, black wordmark.
# Pitch: Probabilistic project intelligence for pharma congress event planning.

def _op_gradient(slide, x, y, w, h):
    base = add_rect(slide, x, y, w, h, OP_DEEP)
    spPr = base.fill._xPr.find(qn("a:solidFill"))
    if spPr is not None:
        base.fill._xPr.remove(spPr)
    g = etree.SubElement(base.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in [(0, OP_DEEP), (50000, OP_ORANGE), (100000, OP_YELLOW)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "2700000"); lin.set("scaled", "1")
    return base


def _project_header(slide, num, project_name, title, *, accent, deep, line):
    """Generic header block reused by Operativo and GSK slides."""
    add_rect(slide, 0, 0, SW, Inches(0.18), accent)
    chip = add_rounded(slide, Inches(0.5), Inches(0.42), Inches(2.4), Inches(0.45),
                       deep, corner=0.3)
    add_text(slide, Inches(0.5), Inches(0.42), Inches(2.4), Inches(0.45),
             project_name, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(3.05), Inches(0.38), Inches(10), Inches(0.55),
             title, size=22, bold=True, color=deep, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.5), Inches(1.0), Inches(12.3), Emu(9525), line)


def _project_footer(slide, page, total, accent, line):
    add_rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), line)
    add_text(slide, Inches(0.5), SH - Inches(0.30), Inches(8), Inches(0.28),
             "Portfolio Pitch Deck  ·  Multi-client AI accelerators",
             size=9, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(2), SH - Inches(0.30), Inches(1.5), Inches(0.28),
             f"{page} / {total}", size=9, color=INK2,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def slide_operativo_intro(idx, total):
    s = prs.slides.add_slide(BLANK)
    _project_header(s, "10", "OPERATIVO  ·  CONGRESS TWIN",
                    "Probabilistic Project Intelligence for Pharma Congress Planning",
                    accent=OP_ORANGE, deep=OP_DEEP, line=OP_LINE)

    # Logo block top-right
    if OPERATIVO_LOGO.exists():
        s.shapes.add_picture(str(OPERATIVO_LOGO),
                             SW - Inches(2.3), Inches(0.45),
                             height=Inches(0.55))

    # Sub-line under header
    add_text(s, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.4),
             "Microsoft Planner sync + Monte Carlo + Markov chain task simulator + AI suggestion engine, "
             "purpose-built for the 18-24 month congress / event planning cycle.",
             size=12, color=INK2)

    # Three cards: Problem / Product / Buyer
    cards = [
        ("The Problem",
         "Pharma medical-affairs congresses are 18-24 month, multi-stakeholder, "
         "high-stakes programs. 30%+ of critical milestones slip; visibility is "
         "fragmented across MS Planner, Teams, email, and ad-hoc trackers."),
        ("The Product",
         "A bidirectional bridge over Microsoft Planner that adds Monte Carlo "
         "timeline simulation, Markov-chain task-state modelling, AI-driven "
         "risk scoring, and an attention dashboard that surfaces blockers, "
         "overdue tasks, and critical-path items in real time."),
        ("Target Buyer",
         "Medical Affairs Congress Ops Lead, Event Director, and Program "
         "Manager at Top-20 pharma manufacturers - the team that owns a "
         "calendar of 50-200 congresses a year and lives in MS Planner today."),
    ]
    cw = (SW - Inches(1.4)) / 3
    cy = Inches(1.7)
    ch = Inches(3.4)
    for i, (h, body) in enumerate(cards):
        cx = Inches(0.5) + i * (cw + Inches(0.2))
        add_rounded(s, cx, cy, cw, ch, WHITE, line=OP_LINE, corner=0.05)
        # left accent in green (matches the logo's green dot)
        add_rect(s, cx, cy, Inches(0.08), ch, OP_GREEN)
        add_text(s, cx + Inches(0.25), cy + Inches(0.22), cw - Inches(0.4),
                 Inches(0.5), h, size=15, bold=True, color=OP_DEEP)
        add_text(s, cx + Inches(0.25), cy + Inches(0.85), cw - Inches(0.4),
                 ch - Inches(1.1), body, size=12, color=INK2)

    # Tech stack ribbon
    tech_y = Inches(5.4)
    add_text(s, Inches(0.5), tech_y, Inches(12.3), Inches(0.3),
             "TECH STACK",
             size=10, bold=True, color=OP_DEEP)
    chips = [
        "MS Graph API (Planner / Teams)", "Monte Carlo + Markov-chain engine",
        "Knowledge Graph semantic layer", "Malloy", "Python + FastAPI",
        "Next.js front-end", "AWS Bedrock-ready agent runtime",
    ]
    x = Inches(0.5); y = tech_y + Inches(0.35); h = Inches(0.32)
    for c in chips:
        approx_w = max(1.6, 0.13 * len(c))
        w = Inches(approx_w)
        add_rounded(s, x, y, w, h, OP_SOFT, line=OP_LINE, corner=0.3)
        add_text(s, x, y, w, h, c, size=10, color=OP_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + w + Inches(0.08)

    # Bottom callout
    cy2 = SH - Inches(1.2)
    ch2 = Inches(0.85)
    cal = add_rounded(s, Inches(0.5), cy2, SW - Inches(1), ch2, OP_DEEP, corner=0.04)
    sf = cal.fill._xPr.find(qn("a:solidFill"))
    cal.fill._xPr.remove(sf)
    g = etree.SubElement(cal.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in [(0, OP_DEEP), (50000, OP_ORANGE), (100000, OP_YELLOW)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "0"); lin.set("scaled", "1")
    add_text(s, Inches(0.8), cy2 + Inches(0.12), SW - Inches(1.6), Inches(0.4),
             "What MS Planner gives you is a task list. Congress Twin gives you a probability distribution.",
             size=15, bold=True, color=WHITE)
    add_text(s, Inches(0.8), cy2 + Inches(0.45), SW - Inches(1.6), Inches(0.35),
             "Stop reacting to overdue items. Start managing the congress against its forecast.",
             size=11, color=WHITE)

    _project_footer(s, idx, total, OP_ORANGE, OP_LINE)


def slide_operativo_value(idx, total):
    s = prs.slides.add_slide(BLANK)
    _project_header(s, "11", "OPERATIVO  ·  CONGRESS TWIN",
                    "What We Bring to Their Table",
                    accent=OP_ORANGE, deep=OP_DEEP, line=OP_LINE)
    add_text(s, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.4),
             "Three concrete outcomes a Congress Ops lead can put on a quarterly board slide within 90 days of go-live.",
             size=12, color=INK2)

    # Three big outcome cards with example mechanics
    outcomes = [
        ("40%", "Reduction in planning cycle time",
         "Example: a 22-month congress program with 240 tasks. The AI suggestion engine "
         "flagged 18 reassignment opportunities (over-loaded owners, mis-routed approvals) "
         "and Monte Carlo simulation re-baselined the timeline within an hour - work that "
         "previously took a planner two weeks of stakeholder calls."),
        ("5,000", "Monte Carlo simulations per run",
         "Example: a congress on the critical path to a regulatory filing. Probabilistic "
         "completion windows let the team commit to P85 dates with confidence instead of "
         "single-point estimates - and triggered automated escalation when P50 slipped "
         "by more than 3 days."),
        ("Real-time", "Risk scoring on every task",
         "Example: a regulator-facing deliverable scored 87/100 risk on Friday afternoon "
         "based on dependency state + assignee workload + historical pattern. The system "
         "proposed two alternative owners; the Ops Lead approved the reassignment in one "
         "click and the score fell to 31/100 the same day."),
    ]
    cw = (SW - Inches(1.4)) / 3
    cy = Inches(1.7)
    ch = Inches(3.7)
    for i, (big, lbl, body) in enumerate(outcomes):
        cx = Inches(0.5) + i * (cw + Inches(0.2))
        add_rounded(s, cx, cy, cw, ch, WHITE, line=OP_LINE, corner=0.05)
        add_rect(s, cx, cy, cw, Inches(0.08), OP_ORANGE)
        add_text(s, cx + Inches(0.25), cy + Inches(0.25), cw - Inches(0.4),
                 Inches(0.85), big, size=40, bold=True, color=OP_ORANGE)
        add_text(s, cx + Inches(0.25), cy + Inches(1.12), cw - Inches(0.4),
                 Inches(0.4), lbl, size=12, bold=True, color=OP_DEEP)
        add_text(s, cx + Inches(0.25), cy + Inches(1.6), cw - Inches(0.4),
                 ch - Inches(1.8), body, size=11, color=INK2)

    # Engagement bar
    eng_y = SH - Inches(1.45)
    eng_h = Inches(1.05)
    add_rounded(s, Inches(0.5), eng_y, SW - Inches(1), eng_h, OP_SOFT, line=OP_LINE, corner=0.04)
    add_text(s, Inches(0.75), eng_y + Inches(0.12), SW - Inches(1.5), Inches(0.35),
             "ENGAGEMENT  ·  8-WEEK PROOF-OF-VALUE",
             size=11, bold=True, color=OP_DEEP)
    add_text(s, Inches(0.75), eng_y + Inches(0.45), SW - Inches(1.5), Inches(0.55),
             "Weeks 1-2: connect to one live congress in MS Planner.  "
             "Weeks 3-6: run Monte Carlo + risk-scoring against historical "
             "data.  Weeks 7-8: deliver attention dashboard and a "
             "side-by-side comparison against the planner-only baseline.",
             size=11, color=INK2)

    _project_footer(s, idx, total, OP_ORANGE, OP_LINE)


# ============================================================================
# GSK VELOCITY  (slides 12-13)
# ============================================================================
# Brand: GSK orange-red gradient (sunset / orange / deep), charcoal text.
# Pitch: AI-led compression of the drug lifecycle, inspired by Novo Nordisk.

def slide_gsk_intro(idx, total):
    s = prs.slides.add_slide(BLANK)
    _project_header(s, "12", "GSK VELOCITY",
                    "AI-Led Compression of the Drug Lifecycle",
                    accent=GSK_ORANGE, deep=GSK_CHARCOAL, line=GSK_LINE)

    if GSK_LOGO.exists():
        s.shapes.add_picture(str(GSK_LOGO),
                             SW - Inches(1.8), Inches(0.32),
                             height=Inches(0.8))

    add_text(s, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.4),
             "A three-module accelerator across Discovery, Clinical, and Pharmacovigilance - "
             "inspired by Novo Nordisk's published compression of LPLV-to-filing cycle time, "
             "oriented around the R&D / Clinical / Safety leadership at GSK.",
             size=12, color=GSK_SLATE)

    cards = [
        ("The Reference Story",
         "Novo Nordisk publicly compressed Last-Patient-Last-Visit to filing materially "
         "with AI applied across data lock, CSR drafting, and submission readiness. GSK "
         "has the appetite, the data, and the talent - this accelerator is the playbook."),
        ("The Product (3 Modules)",
         "DISCOVERY: target identification + de novo molecule design + structure panel.  "
         "CLINICAL: LPLV -> DB Lock -> CSR draft, with cycle-time savings by stage and a "
         "submission-gap auditor.  PHARMACOVIGILANCE: multi-channel ICSR intake, signal "
         "detection map, PSUR co-pilot, literature monitoring."),
        ("Target Buyer",
         "Chief Medical Officer, Head of R&D, Head of Clinical Operations, "
         "Head of Pharmacovigilance at GSK. Cross-functional but anchored in "
         "Clinical as the first vertical."),
    ]
    cw = (SW - Inches(1.4)) / 3
    cy = Inches(1.7)
    ch = Inches(3.5)
    for i, (h, body) in enumerate(cards):
        cx = Inches(0.5) + i * (cw + Inches(0.2))
        add_rounded(s, cx, cy, cw, ch, WHITE, line=GSK_LINE, corner=0.05)
        add_rect(s, cx, cy, Inches(0.08), ch, GSK_ORANGE)
        add_text(s, cx + Inches(0.25), cy + Inches(0.22), cw - Inches(0.4),
                 Inches(0.5), h, size=15, bold=True, color=GSK_CHARCOAL)
        add_text(s, cx + Inches(0.25), cy + Inches(0.85), cw - Inches(0.4),
                 ch - Inches(1.1), body, size=12, color=GSK_SLATE)

    # Module strip with icons
    strip_y = Inches(5.5)
    add_text(s, Inches(0.5), strip_y, Inches(12.3), Inches(0.3),
             "LIFECYCLE COVERAGE",
             size=10, bold=True, color=GSK_CHARCOAL)
    modules = [
        ("Discovery",        "Target ID  ·  De novo design  ·  Structure panel"),
        ("Clinical",         "LPLV  ·  DB Lock  ·  CSR draft  ·  Submission audit"),
        ("Pharmacovigilance","ICSR intake  ·  Signal detection  ·  PSUR  ·  Literature"),
        ("Connectors",       "Internal + external + model-layer integration ecosystem"),
    ]
    mw = (SW - Inches(1.2)) / 4
    my = strip_y + Inches(0.35)
    mh = Inches(0.95)
    for i, (name, sub) in enumerate(modules):
        mx = Inches(0.5) + i * (mw + Inches(0.1))
        # gradient cap
        cap = add_rect(s, mx, my, mw - Inches(0.1), Inches(0.18), GSK_ORANGE)
        sf = cap.fill._xPr.find(qn("a:solidFill"))
        cap.fill._xPr.remove(sf)
        g = etree.SubElement(cap.fill._xPr, qn("a:gradFill"))
        g.set("flip", "none"); g.set("rotWithShape", "1")
        gsLst = etree.SubElement(g, qn("a:gsLst"))
        for pos, col in [(0, GSK_SUNSET), (50000, GSK_ORANGE), (100000, GSK_DEEP)]:
            gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
            sc = etree.SubElement(gs, qn("a:srgbClr"))
            sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
        lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "0"); lin.set("scaled", "1")
        add_rounded(s, mx, my + Inches(0.18), mw - Inches(0.1), mh - Inches(0.18),
                    GSK_CLOUD, line=GSK_LINE, corner=0.04)
        add_text(s, mx + Inches(0.1), my + Inches(0.25), mw - Inches(0.3),
                 Inches(0.3), name, size=12, bold=True, color=GSK_CHARCOAL)
        add_text(s, mx + Inches(0.1), my + Inches(0.58), mw - Inches(0.3),
                 mh - Inches(0.65), sub, size=9, color=GSK_SLATE)

    _project_footer(s, idx, total, GSK_ORANGE, GSK_LINE)


def slide_gsk_value(idx, total):
    s = prs.slides.add_slide(BLANK)
    _project_header(s, "13", "GSK VELOCITY",
                    "What We Bring to Their Table",
                    accent=GSK_ORANGE, deep=GSK_CHARCOAL, line=GSK_LINE)
    add_text(s, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.4),
             "Three measurable wins a GSK R&D leader can land within a single year, with concrete worked examples.",
             size=12, color=GSK_SLATE)

    outcomes = [
        ("60d -> 18d", "LPLV to DB Lock cycle compression",
         "Example: the Clinical co-pilot ingests query backlog + EDC audit trail + "
         "data-management SOPs and auto-drafts ~70% of the data-cleaning queries. The "
         "DM team approves, edits, and routes - data lock happens in weeks not months. "
         "This is the single biggest cycle-time lever Novo Nordisk publicly cited."),
        ("CSR Co-Pilot", "First draft Clinical Study Report",
         "Example: SAP + TFL outputs + protocol go in; a structured first-draft CSR "
         "(Sections 9-16) comes out, with deviation rationale templates, AE summary "
         "narrative, and submission-gap audit. Medical writers move from blank-page "
         "drafting to editorial review - the highest-value use of their time."),
        ("Early Signals", "Multi-channel pharmacovigilance",
         "Example: spontaneous reports + EHR + literature + social are ingested in "
         "parallel; the signal-detection map flags an emerging adverse event 4-6 weeks "
         "earlier than the traditional periodic review would. The PSUR co-pilot drafts "
         "the regulator-facing narrative the same week."),
    ]
    cw = (SW - Inches(1.4)) / 3
    cy = Inches(1.7)
    ch = Inches(3.8)
    for i, (big, lbl, body) in enumerate(outcomes):
        cx = Inches(0.5) + i * (cw + Inches(0.2))
        add_rounded(s, cx, cy, cw, ch, WHITE, line=GSK_LINE, corner=0.05)
        add_rect(s, cx, cy, cw, Inches(0.08), GSK_ORANGE)
        add_text(s, cx + Inches(0.25), cy + Inches(0.3), cw - Inches(0.4),
                 Inches(0.85), big, size=30, bold=True, color=GSK_DEEP)
        add_text(s, cx + Inches(0.25), cy + Inches(1.18), cw - Inches(0.4),
                 Inches(0.4), lbl, size=12, bold=True, color=GSK_CHARCOAL)
        add_text(s, cx + Inches(0.25), cy + Inches(1.65), cw - Inches(0.4),
                 ch - Inches(1.85), body, size=11, color=GSK_SLATE)

    # Engagement bar
    eng_y = SH - Inches(1.35)
    eng_h = Inches(0.95)
    eng = add_rounded(s, Inches(0.5), eng_y, SW - Inches(1), eng_h, GSK_CHARCOAL, corner=0.04)
    sf = eng.fill._xPr.find(qn("a:solidFill"))
    eng.fill._xPr.remove(sf)
    g = etree.SubElement(eng.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in [(0, GSK_CHARCOAL), (50000, GSK_ORANGE), (100000, GSK_DEEP)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "0"); lin.set("scaled", "1")
    add_text(s, Inches(0.75), eng_y + Inches(0.12), SW - Inches(1.5), Inches(0.35),
             "ENGAGEMENT  ·  12-WEEK CLINICAL-FIRST VERTICAL PILOT",
             size=11, bold=True, color=WHITE)
    add_text(s, Inches(0.75), eng_y + Inches(0.45), SW - Inches(1.5), Inches(0.45),
             "Land on Clinical (LPLV -> DB Lock -> CSR), prove cycle-time savings on one therapy area, "
             "then expand to Discovery and Pharmacovigilance in subsequent quarters.",
             size=11, color=WHITE)

    _project_footer(s, idx, total, GSK_ORANGE, GSK_LINE)


# ============================================================================
# BUILD
# ============================================================================

def build():
    TOTAL = 13
    slide_title()                          # 1
    slide_thesis(2, TOTAL)                 # 2
    slide_asset(3, TOTAL)                  # 3
    slide_mapping(4, TOTAL)                # 4
    slide_usecases_a(5, TOTAL)             # 5
    slide_usecases_b(6, TOTAL)             # 6
    slide_position(7, TOTAL)               # 7
    slide_engagement(8, TOTAL)             # 8
    slide_next(9, TOTAL)                   # 9
    slide_operativo_intro(10, TOTAL)       # 10 — Operativo / Congress Twin
    slide_operativo_value(11, TOTAL)       # 11
    slide_gsk_intro(12, TOTAL)             # 12 — GSK Velocity
    slide_gsk_value(13, TOTAL)             # 13

    try:
        prs.save(str(OUT))
        print(f"Wrote {OUT}")
    except PermissionError:
        side = OUT.with_name(OUT.stem + "_v2" + OUT.suffix)
        prs.save(str(side))
        print(f"[note] {OUT.name} is open in PowerPoint - wrote to {side.name} instead.")
        print(f"Close the original and run again, or open {side} directly.")


if __name__ == "__main__":
    build()
