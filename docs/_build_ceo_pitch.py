"""
3-slide CEO pitch for Circana — Circana brand theme, plain neutral language,
Circana logo on every slide.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Circana brand palette (matches /frontend and the long-form deck) ────────
MAGENTA = RGBColor(0xC8, 0x38, 0x9B)
PURPLE  = RGBColor(0x6B, 0x2C, 0x8E)
DEEP    = RGBColor(0x2A, 0x0F, 0x40)
INK     = RGBColor(0x15, 0x15, 0x1A)
INK2    = RGBColor(0x4A, 0x4A, 0x55)
LINE    = RGBColor(0xE6, 0xE1, 0xEE)
SOFT    = RGBColor(0xF7, 0xF4, 0xFB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

ROOT = Path(__file__).resolve().parent.parent
LOGO = ROOT / "circana-logo.png"
OUT  = ROOT / "docs" / "circana_ceo_pitch.pptx"

# Widescreen 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Drawing helpers ──────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill, line=None, corner=0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = corner
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def gradient(shape, stops, angle="2700000"):
    sf = shape.fill._xPr.find(qn("a:solidFill"))
    if sf is not None:
        shape.fill._xPr.remove(sf)
    g = etree.SubElement(shape.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in stops:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", angle); lin.set("scaled", "1")


def corner_logo(slide, on_dark: bool = False):
    """Small Circana logo top-right of every interior slide."""
    if not LOGO.exists():
        return
    plate_w = Inches(1.55)
    plate_h = Inches(0.55)
    plate_x = SW - plate_w - Inches(0.45)
    plate_y = Inches(0.32)
    if on_dark:
        add_rounded(slide, plate_x, plate_y, plate_w, plate_h, WHITE, corner=0.25)
    slide.shapes.add_picture(str(LOGO),
                             plate_x + Inches(0.18),
                             plate_y + Inches(0.10),
                             height=Inches(0.36))


def footer(slide, num):
    add_rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), SOFT)
    add_text(slide, Inches(0.5), SH - Inches(0.30), Inches(9), Inches(0.28),
             "Ask Liquid Data  ·  CEO Brief",
             size=9, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(2), SH - Inches(0.30), Inches(1.5), Inches(0.28),
             f"{num} / 3", size=9, color=INK2,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Where the opportunity is
# ═══════════════════════════════════════════════════════════════════════════
def slide_opportunity():
    s = prs.slides.add_slide(BLANK)
    band = add_rect(s, 0, 0, SW, SH, DEEP)
    gradient(band, [(0, DEEP), (55000, PURPLE), (100000, MAGENTA)], angle="2700000")

    # Decorative circle
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              SW - Inches(3.5), -Inches(2.5),
                              Inches(6), Inches(6))
    deco.fill.solid(); deco.fill.fore_color.rgb = WHITE
    deco.line.fill.background()
    solidFill = deco.fill._xPr.find(qn("a:solidFill"))
    srgb = solidFill.find(qn("a:srgbClr"))
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", "10000")

    # Logo plate top-left (the hero placement on slide 1)
    add_rounded(s, Inches(0.6), Inches(0.5), Inches(2.2), Inches(0.85), WHITE, corner=0.18)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(0.85), Inches(0.6),
                             height=Inches(0.65))

    # Tag top-right
    add_text(s, SW - Inches(5.6), Inches(0.7), Inches(5), Inches(0.4),
             "CEO BRIEF   ·   MAY 2026",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # Hero headline — plain, neutral
    add_text(s, Inches(0.6), Inches(1.9), Inches(11.5), Inches(1.2),
             "Faster answers for your clients —",
             size=42, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(2.75), Inches(11.5), Inches(0.9),
             "without changing how they work today.",
             size=30, bold=False, color=WHITE)

    # Sub thesis — plain language, no jargon
    add_text(s, Inches(0.6), Inches(3.9), Inches(11.5), Inches(1.6),
             ("Circana has the most complete view of CPG and General Merchandise "
              "sales in the market. Today, getting an answer out of that data still "
              "takes a brand team hours or days. We'd like to help you close that "
              "gap with an AI layer that lets people ask questions in plain English "
              "and get a chart, a written answer, and a suggested next step in seconds."),
             size=15, color=WHITE)

    # Hero stat strip
    stats = [
        ("$5.8T",   "Consumer spend you track"),
        ("42M",     "CPG & GM items measured"),
        ("477K",    "Stores in the panel"),
        ("Hours",   "How long an answer takes today"),
    ]
    x0 = Inches(0.6); w = (SW - Inches(1.2)) / 4
    y = Inches(5.85); h = Inches(1.05)
    for i, (n, lbl) in enumerate(stats):
        bx = x0 + i * w + Inches(0.05)
        bw = w - Inches(0.1)
        card = add_rounded(s, bx, y, bw, h, DEEP, corner=0.1)
        sf = card.fill._xPr.find(qn("a:solidFill"))
        srgb = sf.find(qn("a:srgbClr"))
        alpha = etree.SubElement(srgb, qn("a:alpha"))
        alpha.set("val", "35000")
        add_text(s, bx, y + Inches(0.15), bw, Inches(0.45),
                 n, size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bx + Inches(0.1), y + Inches(0.6), bw - Inches(0.2),
                 Inches(0.4), lbl, size=10, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What we'd build with you
# ═══════════════════════════════════════════════════════════════════════════
def slide_product():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, Inches(0.18), MAGENTA)

    chip = add_rounded(s, Inches(0.5), Inches(0.42), Inches(0.7), Inches(0.45),
                       PURPLE, corner=0.3)
    add_text(s, Inches(0.5), Inches(0.42), Inches(0.7), Inches(0.45),
             "02", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.35), Inches(0.38), Inches(8.5), Inches(0.55),
             "What we'd build with you",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    corner_logo(s, on_dark=False)
    add_rect(s, Inches(0.5), Inches(1.0), Inches(12.3), Emu(9525), LINE)
    add_text(s, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.4),
             "Three things that work together. Each is useful on its own.",
             size=12, color=INK2)

    # Left column: three plain-language pieces
    pieces = [
        ("Ask in plain English",
         "A brand or category manager types a question the way they'd ask a colleague: "
         "\"Which brands are losing share in Energy Drinks?\" The system understands the "
         "question, pulls the data, picks the right chart, and writes a 2-sentence answer."),
        ("Set the rules that matter to you",
         "You and your clients define what 'good' looks like: distribution thresholds, "
         "joint business plan targets, price guardrails. The system watches the data "
         "and tells the right person the moment something needs attention."),
        ("Let it do the routine work",
         "Small AI helpers that draft monthly business reviews, look into share changes, "
         "spot new launches that are about to take off, and prepare retailer meeting "
         "briefs. A person always reviews and approves before anything goes out."),
    ]
    lx = Inches(0.5); lw = Inches(5.6)
    ly = Inches(1.7); lh = Inches(1.65)
    for i, (h, body) in enumerate(pieces):
        y = ly + i * (lh + Inches(0.1))
        add_rect(s, lx, y, Inches(0.08), lh, MAGENTA)
        add_text(s, lx + Inches(0.2), y + Inches(0.08), lw - Inches(0.3),
                 Inches(0.35), h, size=13, bold=True, color=DEEP)
        add_text(s, lx + Inches(0.2), y + Inches(0.45), lw - Inches(0.3),
                 lh - Inches(0.5), body, size=11, color=INK2)

    # Right column: 5 things you can ask it to do
    use_cases = [
        ("1", "Ask Liquid Data",
              "Ask anything about share, sales velocity, or distribution."),
        ("2", "Draft Monthly Reviews",
              "Auto-prepares Monthly Business Reviews — analyst edits."),
        ("3", "Why is share moving?",
              "Looks at the data and explains what's driving a gain or loss."),
        ("4", "Spot winning launches early",
              "Flags new products that look like next year's pacesetters."),
        ("5", "Prepare retailer meetings",
              "Pulls together one brief for any top-to-top, 90 days out."),
    ]
    rx = Inches(6.4); rw = Inches(6.4)
    ry = Inches(1.7); rh = Inches(0.86)
    add_text(s, rx, ry - Inches(0.32), rw, Inches(0.3),
             "FIVE THINGS YOU CAN ASK IT TO DO",
             size=10, bold=True, color=PURPLE)
    for i, (num, name, sub) in enumerate(use_cases):
        y = ry + i * (rh + Inches(0.08))
        add_rounded(s, rx, y, rw, rh, WHITE, line=LINE, corner=0.05)
        add_rounded(s, rx + Inches(0.15), y + Inches(0.16),
                    Inches(0.55), Inches(0.55), PURPLE, corner=0.3)
        add_text(s, rx + Inches(0.15), y + Inches(0.16),
                 Inches(0.55), Inches(0.55), num,
                 size=17, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(0.85), y + Inches(0.14),
                 rw - Inches(1.0), Inches(0.32), name,
                 size=13, bold=True, color=DEEP)
        add_text(s, rx + Inches(0.85), y + Inches(0.46),
                 rw - Inches(1.0), Inches(0.4), sub,
                 size=11, color=INK2)

    # Bottom — three plain notes about how it fits
    ny = SH - Inches(1.05)
    add_text(s, Inches(0.5), ny, Inches(12.3), Inches(0.3),
             "HOW IT FITS",
             size=10, bold=True, color=PURPLE)
    notes = [
        "Works with the data you already have — no new systems for your clients to learn.",
        "Works with the AI model of your choice. Easy to swap as the market moves.",
        "Every answer shows where the numbers came from, so anyone can check them.",
    ]
    nx = Inches(0.5); y = ny + Inches(0.35); h = Inches(0.34)
    pw = (SW - Inches(1.0) - Inches(0.10) * 2) / 3
    for c in notes:
        add_rounded(s, nx, y, pw, h, SOFT, line=LINE, corner=0.2)
        add_text(s, nx + Inches(0.15), y, pw - Inches(0.3), h, c,
                 size=10, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
        nx = nx + pw + Inches(0.10)

    footer(s, 2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The benefits & the ask
# ═══════════════════════════════════════════════════════════════════════════
def slide_ask():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, Inches(0.18), MAGENTA)
    chip = add_rounded(s, Inches(0.5), Inches(0.42), Inches(0.7), Inches(0.45),
                       PURPLE, corner=0.3)
    add_text(s, Inches(0.5), Inches(0.42), Inches(0.7), Inches(0.45),
             "03", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.35), Inches(0.38), Inches(8.5), Inches(0.55),
             "The benefits, and what we'd like to do next",
             size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    corner_logo(s, on_dark=False)
    add_rect(s, Inches(0.5), Inches(1.0), Inches(12.3), Emu(9525), LINE)
    add_text(s, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.4),
             "Five practical benefits. A short pilot to prove them on real client work.",
             size=12, color=INK2)

    # 5 benefit tiles — plain labels, no jargon
    benefits = [
        ("5 to 10×",        "Analyst output per day",
         "Your analysts spend less time pulling data and more time thinking about it."),
        ("Same-day",        "Monthly business reviews",
         "Drafted in minutes; the analyst edits instead of starting from scratch."),
        ("2 hours",         "Time to explain a share change",
         "Down from a week or two of cross-team work. Each finding shows its evidence."),
        ("6 months earlier","Spotting pacesetter launches",
         "Continuous tracking surfaces winners well before the annual report would."),
        ("1 day",           "Retailer meeting prep",
         "One brief, current to last week's data. Currently a 5-day team effort."),
    ]
    n = len(benefits)
    tw = (SW - Inches(1.0) - Inches(0.10) * (n - 1)) / n
    ty = Inches(1.7); th = Inches(2.4)
    for i, (big, lbl, sub) in enumerate(benefits):
        tx = Inches(0.5) + i * (tw + Inches(0.10))
        card = add_rounded(s, tx, ty, tw, th, WHITE, line=LINE, corner=0.05)
        cap = add_rect(s, tx, ty, tw, Inches(0.10), MAGENTA)
        gradient(cap, [(0, PURPLE), (100000, MAGENTA)], angle="0")
        # Tile font sized down a bit so longer phrases fit
        big_size = 22 if len(big) <= 8 else 18
        add_text(s, tx + Inches(0.15), ty + Inches(0.3),
                 tw - Inches(0.3), Inches(0.7), big,
                 size=big_size, bold=True, color=MAGENTA, align=PP_ALIGN.CENTER)
        add_text(s, tx + Inches(0.15), ty + Inches(1.05),
                 tw - Inches(0.3), Inches(0.4), lbl,
                 size=11, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        add_text(s, tx + Inches(0.15), ty + Inches(1.5),
                 tw - Inches(0.3), th - Inches(1.6), sub,
                 size=10, color=INK2, align=PP_ALIGN.CENTER)

    # 12-week phases — plain words
    eng_y = Inches(4.4); eng_h = Inches(1.05)
    add_text(s, Inches(0.5), eng_y - Inches(0.32), Inches(12.3), Inches(0.3),
             "12 WEEKS TO PROVE IT  ·  LAND, MEASURE, GROW",
             size=10, bold=True, color=PURPLE)
    phases = [
        ("Weeks 0-2",  "Get aligned",
         "We pick the first use case with you and agree the data we'll use. Two client teams sign up to try it."),
        ("Weeks 3-8",  "Build",
         "We build a working version on real client data, branded as a Circana surface."),
        ("Weeks 9-12", "Try it",
         "The two client teams use it for real work. We measure the time saved and what they'd want next."),
        ("After",      "Grow",
         "Add the other use cases. Make it available to more of your clients as a tier of Liquid Data."),
    ]
    pw = (SW - Inches(1.0) - Inches(0.10) * 3) / 4
    for i, (label, sub, body) in enumerate(phases):
        px = Inches(0.5) + i * (pw + Inches(0.10))
        add_rounded(s, px, eng_y, pw, eng_h, SOFT, line=LINE, corner=0.04)
        add_rect(s, px, eng_y, Inches(0.06), eng_h, MAGENTA)
        add_text(s, px + Inches(0.15), eng_y + Inches(0.08),
                 pw - Inches(0.3), Inches(0.25), label,
                 size=10, bold=True, color=MAGENTA)
        add_text(s, px + Inches(0.15), eng_y + Inches(0.32),
                 pw - Inches(0.3), Inches(0.3), sub,
                 size=12, bold=True, color=DEEP)
        add_text(s, px + Inches(0.15), eng_y + Inches(0.62),
                 pw - Inches(0.3), eng_h - Inches(0.7), body,
                 size=10, color=INK2)

    # Closing — plain ask
    cy = Inches(5.75); ch = Inches(1.25)
    cal = add_rounded(s, Inches(0.5), cy, SW - Inches(1), ch, DEEP, corner=0.04)
    gradient(cal, [(0, DEEP), (55000, PURPLE), (100000, MAGENTA)], angle="0")
    add_text(s, Inches(0.85), cy + Inches(0.18), SW - Inches(1.7), Inches(0.4),
             "WHAT WE'RE ASKING FOR",
             size=10, bold=True, color=WHITE)
    add_text(s, Inches(0.85), cy + Inches(0.45), SW - Inches(1.7), Inches(0.5),
             "A 12-week pilot with two of your clients trying it on real work.",
             size=17, bold=True, color=WHITE)
    add_text(s, Inches(0.85), cy + Inches(0.85), SW - Inches(1.7), Inches(0.35),
             "If it saves the time we think it will, we grow it together. If it doesn't, you've still learned where the gap is.",
             size=12, color=WHITE)

    footer(s, 3)


# ── BUILD ────────────────────────────────────────────────────────────────────
def build():
    slide_opportunity()
    slide_product()
    slide_ask()
    try:
        prs.save(str(OUT))
        print(f"Wrote {OUT}")
    except PermissionError:
        side = OUT.with_name(OUT.stem + "_v2" + OUT.suffix)
        prs.save(str(side))
        print(f"[note] {OUT.name} is open in PowerPoint - wrote to {side.name} instead.")


if __name__ == "__main__":
    build()
