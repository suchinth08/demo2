"""
Auto-MBR Agent — generates a Monthly Business Review deck for a CPG/GM brand
or category, with Circana branding. Orchestrates six data-gather steps, builds
an LLM-narrated insight per section, and renders a downloadable .pptx.

Sections in the generated deck:
  1.  Title
  2.  Executive Summary  (LLM)
  3.  Share Performance
  4.  Distribution & ACV
  5.  Velocity & Pricing
  6.  Promotion Effectiveness
  7.  Panel & Innovation
  8.  Recommended Actions  (LLM)
"""
from __future__ import annotations

import json
import uuid
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from chatbot.services.intent_engine import _call_groq
from chatbot.services.query_engine import get_conn


# ── Storage for generated decks ──────────────────────────────────────────────
OUTPUT_DIR = Path(__file__).parent.parent.parent / "data" / "mbr_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

LOGO_PATH = Path(__file__).parent.parent.parent / "circana-favicon.png"

# ── Brand palette (matches the rest of the project) ──────────────────────────
MAGENTA = RGBColor(0xC8, 0x38, 0x9B)
PURPLE  = RGBColor(0x6B, 0x2C, 0x8E)
DEEP    = RGBColor(0x2A, 0x0F, 0x40)
INK     = RGBColor(0x15, 0x15, 0x1A)
INK2    = RGBColor(0x4A, 0x4A, 0x55)
LINE    = RGBColor(0xE6, 0xE1, 0xEE)
SOFT    = RGBColor(0xF7, 0xF4, 0xFB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)


# ═════════════════════════════════════════════════════════════════════════════
# DATA GATHERING — one function per section
# ═════════════════════════════════════════════════════════════════════════════

def _q(s: str) -> str:
    return s.replace("'", "''")


def _parse_brand_list(brand: Optional[str]) -> list[str]:
    """Split 'Celsius, Coca-Cola' / 'Celsius and Coca-Cola' / 'Celsius' into ['Celsius', 'Coca-Cola']."""
    if not brand:
        return []
    # Accept comma OR ' and ' as separator
    raw = brand.replace(" and ", ",").split(",")
    return [b.strip() for b in raw if b.strip()]


def _scope_filter(brand: Optional[str], category: Optional[str], alias_b="b", alias_c="c") -> str:
    parts = []
    brands = _parse_brand_list(brand)
    if len(brands) == 1:
        parts.append(f"{alias_b}.brand_name = '{_q(brands[0])}'")
    elif len(brands) > 1:
        in_list = ", ".join(f"'{_q(b)}'" for b in brands)
        parts.append(f"{alias_b}.brand_name IN ({in_list})")
    if category:
        parts.append(f"{alias_c}.category_name = '{_q(category)}'")
    return (" AND " + " AND ".join(parts)) if parts else ""


# ── Scope validation: refuse to generate a deck for an empty cell ────────────

def _scope_rowcount(brand: Optional[str], category: Optional[str]) -> int:
    """Cheap pre-flight COUNT(*) on pos_weekly for the requested scope."""
    conn = get_conn()
    sql = f"""
        SELECT COUNT(*) AS n
        FROM pos_weekly p
        JOIN products pr ON p.sku_id = pr.sku_id
        JOIN brands b    ON pr.brand_id = b.brand_id
        JOIN categories c ON pr.category_id = c.category_id
        WHERE CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 26 WEEK
              {_scope_filter(brand, category)}
    """
    rows = conn.execute(sql).fetchdf().to_dict("records")
    return int(rows[0]["n"]) if rows else 0


def _categories_for_brand(brand_list: list[str]) -> list[str]:
    if not brand_list:
        return []
    conn = get_conn()
    in_list = ", ".join(f"'{_q(b)}'" for b in brand_list)
    sql = f"""
        SELECT DISTINCT c.category_name
        FROM products pr
        JOIN brands b ON pr.brand_id = b.brand_id
        JOIN categories c ON pr.category_id = c.category_id
        WHERE b.brand_name IN ({in_list})
        ORDER BY c.category_name
    """
    return [r["category_name"] for r in conn.execute(sql).fetchdf().to_dict("records")]


def _brands_for_category(category: str, limit: int = 12) -> list[str]:
    conn = get_conn()
    sql = f"""
        SELECT b.brand_name, SUM(p.dollars) AS dollars
        FROM pos_weekly p
        JOIN products pr ON p.sku_id = pr.sku_id
        JOIN brands b ON pr.brand_id = b.brand_id
        JOIN categories c ON pr.category_id = c.category_id
        WHERE c.category_name = '{_q(category)}'
        GROUP BY b.brand_name
        ORDER BY dollars DESC
        LIMIT {limit}
    """
    return [r["brand_name"] for r in conn.execute(sql).fetchdf().to_dict("records")]


def _all_brands_existing(brand_list: list[str]) -> tuple[list[str], list[str]]:
    """Split brand_list into (recognized, unknown)."""
    if not brand_list:
        return [], []
    conn = get_conn()
    in_list = ", ".join(f"'{_q(b)}'" for b in brand_list)
    sql = f"SELECT DISTINCT brand_name FROM brands WHERE brand_name IN ({in_list})"
    found = {r["brand_name"] for r in conn.execute(sql).fetchdf().to_dict("records")}
    recognized = [b for b in brand_list if b in found]
    unknown    = [b for b in brand_list if b not in found]
    return recognized, unknown


class ScopeValidationError(ValueError):
    """Raised when the requested MBR scope has no data — carries suggestions."""
    def __init__(self, message: str, suggestions: dict):
        super().__init__(message)
        self.suggestions = suggestions


def _validate_scope(brand: Optional[str], category: Optional[str]) -> None:
    """Raise ScopeValidationError with helpful suggestions if the scope is empty."""
    brand_list = _parse_brand_list(brand)

    # 1. Unknown brand names
    if brand_list:
        recognized, unknown = _all_brands_existing(brand_list)
        if unknown and not recognized:
            raise ScopeValidationError(
                f"None of these brands exist in the dataset: {', '.join(unknown)}.",
                {"unknown_brands": unknown,
                 "hint": "Try one of the demo brands: Liquid Death, Celsius, Coca-Cola, iPhone, LEGO."},
            )
        if unknown:
            # Some recognized, some not — proceed with the recognized ones (warning only)
            pass

    # 2. Scope has data?
    n = _scope_rowcount(brand, category)
    if n > 0:
        return

    # 3. Build helpful suggestions
    suggestions: dict = {}
    msg_parts: list[str] = []
    if brand_list and category:
        valid_cats = _categories_for_brand(brand_list)
        suggestions["valid_categories_for_brand"] = valid_cats
        if valid_cats:
            msg_parts.append(
                f"'{', '.join(brand_list)}' is not present in '{category}'. "
                f"It is in: {', '.join(valid_cats)}."
            )
        else:
            msg_parts.append(f"No data for '{', '.join(brand_list)}'.")
    elif category and not brand_list:
        suggestions["top_brands_in_category"] = _brands_for_category(category)
        msg_parts.append(
            f"No data for category '{category}' in the last 26 weeks. "
            f"Top brands available: {', '.join(suggestions['top_brands_in_category'][:5])}."
        )
    elif brand_list and not category:
        valid_cats = _categories_for_brand(brand_list)
        suggestions["valid_categories_for_brand"] = valid_cats
        msg_parts.append(
            f"No recent data for '{', '.join(brand_list)}'. "
            f"Known categories: {', '.join(valid_cats) or 'none'}."
        )
    else:
        msg_parts.append("Provide a brand or a category.")

    raise ScopeValidationError(" ".join(msg_parts), suggestions)


def _gather_share(brand: Optional[str], category: Optional[str]) -> dict:
    conn = get_conn()
    sql = f"""
        WITH brand_period AS (
          SELECT
            CASE
              WHEN CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
                   THEN 'L13W'
              WHEN CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 26 WEEK
                   AND CAST(p.week_ending AS DATE) <  (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
                   THEN 'P13W'
            END AS period_label,
            b.brand_name, c.category_name, pr.category_id,
            SUM(p.dollars) AS dollars
          FROM pos_weekly p
          JOIN products pr ON p.sku_id = pr.sku_id
          JOIN brands b ON pr.brand_id = b.brand_id
          JOIN categories c ON pr.category_id = c.category_id
          WHERE CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 26 WEEK
                {_scope_filter(brand, category)}
          GROUP BY period_label, b.brand_name, c.category_name, pr.category_id
        ),
        cat_period AS (
          SELECT pr.category_id,
            CASE
              WHEN CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
                   THEN 'L13W'
              WHEN CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 26 WEEK
                   AND CAST(p.week_ending AS DATE) <  (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
                   THEN 'P13W'
            END AS period_label,
            SUM(p.dollars) AS cat_dollars
          FROM pos_weekly p
          JOIN products pr ON p.sku_id = pr.sku_id
          WHERE pr.category_id IN (SELECT DISTINCT category_id FROM brand_period)
            AND CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 26 WEEK
          GROUP BY pr.category_id, period_label
        )
        SELECT bp.brand_name, bp.category_name,
          ROUND(SUM(CASE WHEN bp.period_label='L13W' THEN bp.dollars END) / 1e6, 2) AS current_dollars_mm,
          ROUND(SUM(CASE WHEN bp.period_label='P13W' THEN bp.dollars END) / 1e6, 2) AS prior_dollars_mm,
          ROUND(100.0 * SUM(CASE WHEN bp.period_label='L13W' THEN bp.dollars END)
              / NULLIF(SUM(CASE WHEN bp.period_label='L13W' THEN cp.cat_dollars END), 0), 2) AS current_share_pct,
          ROUND(100.0 * SUM(CASE WHEN bp.period_label='P13W' THEN bp.dollars END)
              / NULLIF(SUM(CASE WHEN bp.period_label='P13W' THEN cp.cat_dollars END), 0), 2) AS prior_share_pct
        FROM brand_period bp
        JOIN cat_period cp ON bp.category_id = cp.category_id AND bp.period_label = cp.period_label
        WHERE bp.period_label IS NOT NULL
        GROUP BY bp.brand_name, bp.category_name
        ORDER BY current_share_pct DESC
        LIMIT 10
    """
    rows = conn.execute(sql).fetchdf().to_dict("records")
    for r in rows:
        cs = r.get("current_share_pct") or 0
        ps = r.get("prior_share_pct") or 0
        r["share_pt_change"] = round(cs - ps, 2)
    return {"rows": rows, "title": "Share Performance (L13W vs Prior 13W)"}


def _gather_distribution(brand: Optional[str], category: Optional[str]) -> dict:
    conn = get_conn()
    sql = f"""
        SELECT
          r.retailer_name,
          b.brand_name,
          ROUND(AVG(p.acv_weighted_distribution_pct), 1) AS acv_pct,
          ROUND(100 - AVG(p.acv_weighted_distribution_pct), 1) AS distribution_gap_pct,
          ROUND(SUM(p.dollars)/1e6, 2) AS dollars_mm
        FROM pos_weekly p
        JOIN products pr ON p.sku_id = pr.sku_id
        JOIN brands b ON pr.brand_id = b.brand_id
        JOIN categories c ON pr.category_id = c.category_id
        JOIN retailers r ON p.retailer_id = r.retailer_id
        WHERE CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
              {_scope_filter(brand, category)}
        GROUP BY r.retailer_name, b.brand_name
        ORDER BY dollars_mm DESC
        LIMIT 12
    """
    rows = conn.execute(sql).fetchdf().to_dict("records")
    return {"rows": rows, "title": "Distribution & ACV by Retailer (L13W)"}


def _gather_velocity_pricing(brand: Optional[str], category: Optional[str]) -> dict:
    conn = get_conn()
    sql = f"""
        WITH cat_avg AS (
          SELECT pr.category_id, AVG(p.avg_price) AS cat_avg_price
          FROM pos_weekly p JOIN products pr ON p.sku_id = pr.sku_id
          WHERE CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
          GROUP BY pr.category_id
        )
        SELECT
          b.brand_name, c.category_name,
          ROUND(SUM(p.dollars) / NULLIF(SUM(p.stores_selling), 0), 2) AS dollars_per_store_per_week,
          ROUND(SUM(p.units)   / NULLIF(SUM(p.stores_selling), 0), 2) AS units_per_store_per_week,
          ROUND(AVG(p.avg_price), 2) AS avg_price,
          ROUND(100.0 * AVG(p.avg_price) / MAX(ca.cat_avg_price), 1) AS price_index_vs_category
        FROM pos_weekly p
        JOIN products pr ON p.sku_id = pr.sku_id
        JOIN brands b ON pr.brand_id = b.brand_id
        JOIN categories c ON pr.category_id = c.category_id
        JOIN cat_avg ca ON pr.category_id = ca.category_id
        WHERE CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
              {_scope_filter(brand, category)}
        GROUP BY b.brand_name, c.category_name
        ORDER BY dollars_per_store_per_week DESC
        LIMIT 10
    """
    rows = conn.execute(sql).fetchdf().to_dict("records")
    return {"rows": rows, "title": "Velocity & Pricing (L13W)"}


def _gather_promotion(brand: Optional[str], category: Optional[str]) -> dict:
    conn = get_conn()
    sql = f"""
        SELECT
          b.brand_name, c.category_name,
          ROUND(100.0 * SUM(p.incremental_units) / NULLIF(SUM(p.base_units), 0), 1) AS promo_lift_pct,
          ROUND(100.0 * SUM(CASE WHEN p.on_promo THEN p.units END)
              / NULLIF(SUM(p.units), 0), 1) AS pct_volume_on_promo,
          ROUND(SUM(p.incremental_dollars) / 1e6, 2) AS incremental_dollars_mm
        FROM pos_weekly p
        JOIN products pr ON p.sku_id = pr.sku_id
        JOIN brands b ON pr.brand_id = b.brand_id
        JOIN categories c ON pr.category_id = c.category_id
        WHERE CAST(p.week_ending AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 13 WEEK
              {_scope_filter(brand, category)}
        GROUP BY b.brand_name, c.category_name
        ORDER BY promo_lift_pct DESC
        LIMIT 10
    """
    rows = conn.execute(sql).fetchdf().to_dict("records")
    return {"rows": rows, "title": "Promotion Effectiveness (L13W)"}


def _gather_panel_innovation(brand: Optional[str], category: Optional[str]) -> dict:
    conn = get_conn()
    panel_sql = f"""
        SELECT b.brand_name, c.category_name,
          ROUND(AVG(ph.hh_penetration_pct), 1) AS hh_penetration_pct,
          ROUND(AVG(ph.trips_per_buyer), 2)    AS trips_per_buyer,
          ROUND(AVG(ph.dollars_per_trip), 2)   AS dollars_per_trip,
          ROUND(AVG(ph.repeat_rate_pct), 1)    AS repeat_rate_pct
        FROM panel_household ph
        JOIN brands b ON ph.brand_id = b.brand_id
        JOIN categories c ON b.primary_category_id = c.category_id
        WHERE CAST(ph.period_month AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 6 MONTH
              {_scope_filter(brand, category)}
        GROUP BY b.brand_name, c.category_name
        ORDER BY hh_penetration_pct DESC
        LIMIT 8
    """
    innovation_sql = f"""
        SELECT b.brand_name, c.category_name, pr.launch_date,
          ROUND(SUM(p.dollars)/1e6, 2) AS year1_dollars_mm,
          ROUND(AVG(p.acv_weighted_distribution_pct), 1) AS avg_acv_pct
        FROM pos_weekly p
        JOIN products pr ON p.sku_id = pr.sku_id
        JOIN brands b ON pr.brand_id = b.brand_id
        JOIN categories c ON pr.category_id = c.category_id
        WHERE CAST(pr.launch_date AS DATE) >= (SELECT MAX(CAST(week_ending AS DATE)) FROM pos_weekly) - INTERVAL 18 MONTH
              {_scope_filter(brand, category)}
        GROUP BY b.brand_name, c.category_name, pr.launch_date
        HAVING SUM(p.dollars) > 50000
        ORDER BY year1_dollars_mm DESC
        LIMIT 6
    """
    return {
        "panel": conn.execute(panel_sql).fetchdf().to_dict("records"),
        "innovation": conn.execute(innovation_sql).fetchdf().to_dict("records"),
        "title": "Panel & Innovation",
    }


# ═════════════════════════════════════════════════════════════════════════════
# NARRATION — LLM passes for per-section + executive summary + recommendations
# ═════════════════════════════════════════════════════════════════════════════

def _payload_is_empty(payload: dict) -> bool:
    """True if a gather function returned no usable rows."""
    if not isinstance(payload, dict):
        return True
    # Most sections wrap data in "rows"
    if "rows" in payload:
        return not payload.get("rows")
    # Panel & innovation has two row collections
    if "panel" in payload or "innovation" in payload:
        return not (payload.get("panel") or payload.get("innovation"))
    return True


def _narrate_section(section_name: str, scope: dict, payload: dict) -> str:
    """Produce a 2-3 sentence section commentary grounded in the payload."""
    # Deterministic fallback when there is genuinely no data — never let the
    # LLM invent a "data unavailable" narrative.
    if _payload_is_empty(payload):
        label = scope.get("brand") or scope.get("category") or "this scope"
        return (
            f"No data available for {label} in this section over the last 13 weeks. "
            "This typically means the brand is not present in the requested category "
            "or the scope is too narrow. Try broadening the scope."
        )

    system = (
        "You are a CPG/GM market insights analyst writing a Monthly Business "
        "Review. Write 2-3 concise, professional sentences grounded in the data "
        "rows below. Use industry vocabulary (share, ACV, velocity, lift, "
        "penetration). Never invent numbers that are not in the data. "
        "If the data rows are empty, say so plainly — do not hallucinate."
    )
    prompt = (
        f"Scope: {json.dumps(scope)}\n"
        f"Section: {section_name}\n\n"
        f"Data (top rows):\n{json.dumps(payload, default=str, indent=2)[:2200]}\n\n"
        "Write the 2-3 sentence section commentary."
    )
    try:
        return _call_groq(system, [{"role": "user", "content": prompt}], max_tokens=240)
    except Exception:
        return ""


def _executive_summary(scope: dict, sections: dict) -> str:
    system = (
        "You are a CPG/GM market insights analyst writing the Executive Summary "
        "of a Monthly Business Review. In 3-4 sentences, summarise the most "
        "important finding from each section and surface the single biggest "
        "story the leader needs to know. Use specific numbers from the section "
        "data; never invent numbers."
    )
    prompt = (
        f"Scope: {json.dumps(scope)}\n"
        f"Section narratives:\n" +
        "\n".join(f"- {k}: {v.get('narrative','')}" for k, v in sections.items()) +
        "\n\nWrite the executive summary."
    )
    try:
        return _call_groq(system, [{"role": "user", "content": prompt}], max_tokens=320)
    except Exception:
        return ""


def _recommended_actions(scope: dict, sections: dict) -> str:
    system = (
        "You are a CPG/GM market insights analyst. List 3-5 specific, "
        "actionable recommendations that follow directly from the MBR findings "
        "below. Each recommendation must be one sentence, start with an "
        "imperative verb, and reference a concrete metric or driver."
    )
    prompt = (
        f"Scope: {json.dumps(scope)}\n"
        f"Section narratives:\n" +
        "\n".join(f"- {k}: {v.get('narrative','')}" for k, v in sections.items()) +
        "\n\nList the 3-5 recommendations (one per line, start with '- ')."
    )
    try:
        return _call_groq(system, [{"role": "user", "content": prompt}], max_tokens=380)
    except Exception:
        return ""


# ═════════════════════════════════════════════════════════════════════════════
# PPTX RENDERING
# ═════════════════════════════════════════════════════════════════════════════

def _add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _add_rounded(slide, x, y, w, h, fill, line=None, corner=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = corner
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def _gradient_band(slide, x, y, w, h):
    base = _add_rect(slide, x, y, w, h, DEEP)
    spPr = base.fill._xPr.find(qn("a:solidFill"))
    if spPr is not None:
        base.fill._xPr.remove(spPr)
    gradFill = etree.SubElement(base.fill._xPr, qn("a:gradFill"))
    gradFill.set("flip", "none"); gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, col in [(0, DEEP), (55000, PURPLE), (100000, MAGENTA)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(gradFill, qn("a:lin")); lin.set("ang", "2700000"); lin.set("scaled", "1")
    return base


def _header(slide, sw, sh, section_num: str, title: str, sub: Optional[str] = None):
    _add_rect(slide, 0, 0, sw, Inches(0.18), MAGENTA)
    chip = _add_rounded(slide, Inches(0.5), Inches(0.42), Inches(0.85), Inches(0.45),
                        PURPLE, corner=0.3)
    _add_text(slide, Inches(0.5), Inches(0.42), Inches(0.85), Inches(0.45),
              section_num, size=14, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, Inches(1.5), Inches(0.38), Inches(11.5), Inches(0.55),
              title, size=22, bold=True, color=DEEP, anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, Inches(0.5), Inches(1.0), Inches(12.3), Emu(9525), LINE)
    if sub:
        _add_text(slide, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.4),
                  sub, size=12, color=INK2, italic=True)


def _footer(slide, sw, sh, brand_scope: str, page: int, total: int):
    _add_rect(slide, 0, sh - Inches(0.32), sw, Inches(0.32), SOFT)
    _add_text(slide, Inches(0.5), sh - Inches(0.30), Inches(9), Inches(0.28),
              f"Monthly Business Review  ·  {brand_scope}",
              size=9, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, sw - Inches(2), sh - Inches(0.30), Inches(1.5), Inches(0.28),
              f"{page} / {total}", size=9, color=INK2,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _render_table(slide, rows: list[dict], x, y, w, h,
                  header_cols: Optional[list[str]] = None):
    if not rows:
        _add_text(slide, x, y, w, Inches(0.4), "No data available for this section.",
                  size=12, color=INK2, italic=True)
        return
    cols = header_cols or list(rows[0].keys())
    n_rows = len(rows) + 1  # +header
    n_cols = len(cols)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table

    # Header
    for ci, col in enumerate(cols):
        cell = tbl.cell(0, ci)
        cell.text = _humanize(col)
        cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE
                r.font.name = "Calibri"

    # Body
    for ri, row in enumerate(rows, start=1):
        bg = WHITE if ri % 2 == 1 else SOFT
        for ci, col in enumerate(cols):
            cell = tbl.cell(ri, ci)
            val = row.get(col)
            cell.text = _fmt(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10); r.font.color.rgb = INK
                    r.font.name = "Calibri"


def _humanize(s: str) -> str:
    return s.replace("_pct", " %").replace("_mm", " ($MM)").replace("_", " ").title()


def _fmt(v) -> str:
    if v is None: return "-"
    if isinstance(v, float): return f"{v:,.2f}"
    if isinstance(v, int): return f"{v:,}"
    if isinstance(v, datetime): return v.strftime("%Y-%m-%d")
    s = str(v)
    if "T00:00:00" in s: s = s.split("T")[0]
    return s


def _commentary_card(slide, x, y, w, h, text: str):
    _add_rounded(slide, x, y, w, h, SOFT, line=LINE, corner=0.04)
    _add_rect(slide, x, y, Inches(0.08), h, MAGENTA)
    _add_text(slide, x + Inches(0.2), y + Inches(0.08), w - Inches(0.3),
              Inches(0.3), "ANALYST COMMENTARY",
              size=10, bold=True, color=PURPLE)
    _add_text(slide, x + Inches(0.2), y + Inches(0.42), w - Inches(0.3),
              h - Inches(0.5), text or "(no commentary)",
              size=11, color=INK)


# ── Build the deck ──────────────────────────────────────────────────────────

def _build_pptx(scope_label: str, scope_dict: dict, sections: dict,
                exec_summary: str, recommendations: str, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]

    TOTAL = 8

    # ── Slide 1: Title ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(BLANK)
    _gradient_band(s, 0, 0, SW, SH)
    if LOGO_PATH.exists():
        plate = _add_rounded(s, Inches(0.6), Inches(0.55), Inches(0.9), Inches(0.9),
                             WHITE, corner=0.2)
        s.shapes.add_picture(str(LOGO_PATH), Inches(0.73), Inches(0.68),
                             height=Inches(0.65))
    _add_text(s, Inches(1.7), Inches(0.7), Inches(7), Inches(0.45),
              "MONTHLY BUSINESS REVIEW", size=12, bold=True, color=WHITE)
    _add_text(s, Inches(1.7), Inches(1.05), Inches(7), Inches(0.4),
              "Ask Liquid Data  ·  Market Insights AI", size=11, color=WHITE)

    _add_text(s, Inches(0.6), Inches(2.8), Inches(11.5), Inches(1.4),
              scope_label, size=46, bold=True, color=WHITE)
    _add_text(s, Inches(0.6), Inches(3.9), Inches(11.5), Inches(0.6),
              "Last 13 Weeks  ·  Performance, Distribution, Velocity, Promo, Panel & Innovation",
              size=16, color=WHITE)
    _add_text(s, Inches(0.6), SH - Inches(0.6), Inches(11.5), Inches(0.3),
              f"Generated {datetime.utcnow().strftime('%B %Y')}  ·  Powered by Ask Liquid Data",
              size=11, color=WHITE)

    # ── Slide 2: Executive Summary ──────────────────────────────────────────
    s = prs.slides.add_slide(BLANK)
    _header(s, SW, SH, "01", "Executive Summary")
    _add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5),
              exec_summary or "(no summary generated)",
              size=15, color=INK)
    _footer(s, SW, SH, scope_label, 2, TOTAL)

    # ── Slides 3-7: One per data section ────────────────────────────────────
    section_order = ["share", "distribution", "velocity_pricing", "promotion", "panel_innovation"]
    section_titles = {
        "share":            ("02", "Share Performance",        "Brand share within category, last 13 weeks vs prior 13 weeks"),
        "distribution":     ("03", "Distribution & ACV",        "ACV-weighted distribution and dollar contribution by retailer, L13W"),
        "velocity_pricing": ("04", "Velocity & Pricing",        "Rate of sale and price index vs category, L13W"),
        "promotion":        ("05", "Promotion Effectiveness",   "Lift, % volume on promo, incremental dollars, L13W"),
        "panel_innovation": ("06", "Panel & Innovation",        "Household reach, basket dynamics, recent launches"),
    }
    for page_offset, sec_key in enumerate(section_order, start=3):
        sec = sections[sec_key]
        num, title, sub = section_titles[sec_key]
        s = prs.slides.add_slide(BLANK)
        _header(s, SW, SH, num, title, sub)
        # Layout: table on left (8.5"), commentary card on right (3.6")
        if sec_key == "panel_innovation":
            # Two stacked tables
            panel_rows = sec.get("payload", {}).get("panel", []) or sec.get("panel", [])
            inno_rows  = sec.get("payload", {}).get("innovation", []) or sec.get("innovation", [])
            _add_text(s, Inches(0.5), Inches(1.45), Inches(8.5), Inches(0.3),
                      "Household Panel", size=12, bold=True, color=PURPLE)
            _render_table(s, panel_rows[:5], Inches(0.5), Inches(1.8),
                          Inches(8.5), Inches(2.0))
            _add_text(s, Inches(0.5), Inches(4.0), Inches(8.5), Inches(0.3),
                      "Recent Launches (last 18 months)", size=12, bold=True, color=PURPLE)
            _render_table(s, inno_rows[:5], Inches(0.5), Inches(4.35),
                          Inches(8.5), Inches(2.0))
        else:
            rows = sec.get("payload", {}).get("rows", []) or sec.get("rows", [])
            _render_table(s, rows[:8], Inches(0.5), Inches(1.5),
                          Inches(8.5), Inches(4.8))
        _commentary_card(s, Inches(9.2), Inches(1.5), Inches(3.6), Inches(5.2),
                         sec.get("narrative", ""))
        _footer(s, SW, SH, scope_label, page_offset, TOTAL)

    # ── Slide 8: Recommended Actions ────────────────────────────────────────
    s = prs.slides.add_slide(BLANK)
    _header(s, SW, SH, "07", "Recommended Actions")
    _add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5),
              recommendations or "(no recommendations generated)",
              size=14, color=INK)
    # Closing band
    cy = SH - Inches(1.0)
    cal = _add_rounded(s, Inches(0.5), cy, SW - Inches(1), Inches(0.7), DEEP, corner=0.04)
    sf = cal.fill._xPr.find(qn("a:solidFill"))
    cal.fill._xPr.remove(sf)
    g = etree.SubElement(cal.fill._xPr, qn("a:gradFill"))
    g.set("flip", "none"); g.set("rotWithShape", "1")
    gsLst = etree.SubElement(g, qn("a:gsLst"))
    for pos, col in [(0, DEEP), (100000, MAGENTA)]:
        gs = etree.SubElement(gsLst, qn("a:gs")); gs.set("pos", str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", "{:02X}{:02X}{:02X}".format(col[0], col[1], col[2]))
    lin = etree.SubElement(g, qn("a:lin")); lin.set("ang", "0"); lin.set("scaled", "1")
    _add_text(s, Inches(0.8), cy + Inches(0.15), SW - Inches(1.6), Inches(0.4),
              "Auto-generated by Ask Liquid Data",
              size=14, bold=True, color=WHITE)
    _footer(s, SW, SH, scope_label, 8, TOTAL)

    prs.save(str(out_path))


# ═════════════════════════════════════════════════════════════════════════════
# Public entrypoint
# ═════════════════════════════════════════════════════════════════════════════

def generate_mbr(brand_name: Optional[str] = None,
                 category_name: Optional[str] = None) -> dict:
    """Generate the MBR. Returns metadata + a download token for the .pptx."""
    if not brand_name and not category_name:
        raise ValueError("Provide at least one of brand_name or category_name")
    scope = {"brand": brand_name, "category": category_name}
    scope_label = brand_name or category_name

    # 0. Pre-flight: confirm the scope has data before doing the expensive work
    _validate_scope(brand_name, category_name)

    # 1. Gather data (sequential is fine — DuckDB queries are sub-second each)
    sections = {
        "share":            {"payload": _gather_share(brand_name, category_name)},
        "distribution":     {"payload": _gather_distribution(brand_name, category_name)},
        "velocity_pricing": {"payload": _gather_velocity_pricing(brand_name, category_name)},
        "promotion":        {"payload": _gather_promotion(brand_name, category_name)},
        "panel_innovation": {"payload": _gather_panel_innovation(brand_name, category_name)},
    }

    # 2. Narrate sections in parallel — each LLM call is ~20s through claude CLI,
    #    so 5 sequential = ~100s, parallel = ~20-25s.
    with ThreadPoolExecutor(max_workers=6) as pool:
        narr_futures = {
            k: pool.submit(_narrate_section, k, scope, s["payload"])
            for k, s in sections.items()
        }
        for k, fut in narr_futures.items():
            sections[k]["narrative"] = fut.result()

    # 3. Executive summary + recommendations also in parallel (both depend on
    #    section narratives, but they're independent of each other).
    with ThreadPoolExecutor(max_workers=2) as pool:
        es_fut = pool.submit(_executive_summary, scope, sections)
        rec_fut = pool.submit(_recommended_actions, scope, sections)
        exec_summary = es_fut.result()
        recommendations = rec_fut.result()

    # 4. Render the deck
    token = str(uuid.uuid4())
    safe_label = "".join(c if c.isalnum() else "_" for c in scope_label)[:40]
    filename = f"MBR_{safe_label}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}_{token[:8]}.pptx"
    out_path = OUTPUT_DIR / filename
    _build_pptx(scope_label, scope, sections, exec_summary, recommendations, out_path)

    return {
        "token":            token,
        "filename":         filename,
        "scope_label":      scope_label,
        "brand_name":       brand_name,
        "category_name":    category_name,
        "executive_summary": exec_summary,
        "recommendations":  recommendations,
        "section_narratives": {k: v["narrative"] for k, v in sections.items()},
        "section_data":     {k: v["payload"] for k, v in sections.items()},
        "generated_at":     datetime.utcnow().isoformat(),
    }


def resolve_download_path(filename: str) -> Optional[Path]:
    """Validate the requested filename is a real file in the MBR output dir."""
    safe = Path(filename).name  # strip any path traversal
    p = OUTPUT_DIR / safe
    return p if p.exists() and p.is_file() else None
